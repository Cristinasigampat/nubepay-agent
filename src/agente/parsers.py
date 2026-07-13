"""
parsers.py

Funciones para leer documentos de distintos formatos y devolver siempre
lo mismo: un string con el texto plano del documento.

Uso típico desde otro archivo:
    from agente.parsers import leer_documento
    texto = leer_documento(Path("docs/rh/Politica_de_beneficios_y_vacaciones.pdf"))
"""

from pathlib import Path
from pypdf import PdfReader
from docx import Document
import pandas as pd
from pptx import Presentation
from bs4 import BeautifulSoup
import json
import re


def leer_texto_plano(ruta: Path) -> str:
    """Lee archivos que ya son texto plano: Markdown (.md), texto (.txt)."""
    with open(ruta, "r", encoding="utf-8") as archivo:
        return archivo.read()


def leer_pdf(ruta: Path) -> str:
    """Extrae el texto de todas las páginas de un PDF."""
    lector = PdfReader(ruta)
    texto = ""
    for pagina in lector.pages:
        texto += pagina.extract_text() + "\n"
    return texto


def leer_word(ruta: Path) -> str:
    """Extrae el texto de todos los párrafos de un documento Word."""
    documento = Document(ruta)
    texto = ""
    for parrafo in documento.paragraphs:
        texto += parrafo.text + "\n"
    return texto


def leer_excel(ruta: Path) -> str:
    """
    Convierte una hoja de Excel a texto, UNA LÍNEA POR FILA, repitiendo
    el nombre de cada columna junto a su valor. Esto evita que, al trocear
    la tabla en varios chunks, se pierda de vista a qué columna
    corresponde cada número (algo que SÍ pasaba con to_string()).

    Nota: nuestros archivos de Excel tienen 2 filas de título + 1 fila
    vacía antes del encabezado real de la tabla (fila 4). Por eso usamos
    header=3 (índice 0, o sea la 4ta fila) para que pandas identifique
    bien los nombres de columna reales, en vez de inventar columnas
    "Unnamed: 1", "Unnamed: 2", etc.
    """
    tabla = pd.read_excel(ruta, header=3)
    lineas = []
    for _, fila in tabla.iterrows():
        partes = [f"{columna}: {valor}" for columna, valor in fila.items()]
        lineas.append(" | ".join(partes))
    return "\n".join(lineas)


def leer_csv(ruta: Path) -> str:
    """Convierte un CSV a texto, una línea por fila (mismo criterio que leer_excel)."""
    tabla = pd.read_csv(ruta)
    lineas = []
    for _, fila in tabla.iterrows():
        partes = [f"{columna}: {valor}" for columna, valor in fila.items()]
        lineas.append(" | ".join(partes))
    return "\n".join(lineas)


def leer_powerpoint(ruta: Path) -> str:
    """
    Extrae el texto de todas las cajas de texto de todas las diapositivas,
    y también las notas del orador (que suelen tener contexto adicional
    importante, no visible en la diapositiva en sí).
    """
    presentacion = Presentation(ruta)
    texto = ""
    for diapositiva in presentacion.slides:
        for forma in diapositiva.shapes:
            if forma.has_text_frame:
                texto += forma.text_frame.text + "\n"

        if diapositiva.has_notes_slide:
            notas = diapositiva.notes_slide.notes_text_frame.text
            if notas.strip():  # solo agregamos si la nota no está vacía
                texto += f"[Notas del orador: {notas}]\n"
    return texto


def leer_html(ruta: Path) -> str:
    """Lee un HTML y devuelve solo el texto legible, sin etiquetas."""
    with open(ruta, "r", encoding="utf-8") as archivo:
        contenido = archivo.read()
    sopa = BeautifulSoup(contenido, "html.parser")
    return sopa.get_text()


def leer_json(ruta: Path) -> str:
    """Lee un JSON y lo devuelve como texto legible (no como diccionario)."""
    with open(ruta, "r", encoding="utf-8") as archivo:
        datos = json.load(archivo)
    # json.dumps convierte el diccionario de vuelta a texto, pero prolijo e indentado
    return json.dumps(datos, indent=2, ensure_ascii=False)

def reformatear_tablas_markdown(texto: str) -> str:
    """
    Detecta tablas escritas en formato Markdown (con | y una fila de
    guiones separadora) y las convierte a una línea por fila, con cada
    valor pegado al nombre de su columna -- el mismo criterio que ya
    usamos en leer_excel() y leer_csv(). Esto evita que un LLM confunda
    a qué columna corresponde cada número.
 
    Ejemplo de entrada:
        | Plan | Límite diario | Límite mensual |
        |------|-----------------|-----------------|
        | Básico | USD 150 | USD 1.000 |
 
    Ejemplo de salida:
        Plan: Básico | Límite diario: USD 150 | Límite mensual: USD 1.000
    """
    lineas = texto.split("\n")
    resultado = []
    i = 0
 
    while i < len(lineas):
        linea_actual = lineas[i].strip()
        es_posible_encabezado = linea_actual.startswith("|") and linea_actual.endswith("|")
 
        hay_separador_despues = False
        if es_posible_encabezado and i + 1 < len(lineas):
            siguiente = lineas[i + 1].strip()
            hay_separador_despues = bool(re.match(r"^\|[\s:\-|]+\|$", siguiente))
 
        if es_posible_encabezado and hay_separador_despues:
            # Encontramos una tabla: leemos el encabezado
            columnas = [c.strip() for c in linea_actual.strip("|").split("|")]
            i += 2  # saltamos el encabezado y la fila de guiones
 
            # Procesamos cada fila de datos, mientras sigan empezando con "|"
            while i < len(lineas) and lineas[i].strip().startswith("|"):
                valores = [v.strip() for v in lineas[i].strip().strip("|").split("|")]
                partes = [f"{col}: {val}" for col, val in zip(columnas, valores)]
                resultado.append(" | ".join(partes))
                i += 1
        else:
            # No es parte de una tabla: dejamos la línea tal cual
            resultado.append(lineas[i])
            i += 1
 
    return "\n".join(resultado)


def limpiar_texto(texto: str) -> str:
    """
    Limpia "ruido" común en el texto extraído: espacios/tabs repetidos,
    y más de 2 saltos de línea seguidos (que suelen aparecer como
    residuo de la extracción, sobre todo en PDF).
    """
    texto = re.sub(r"[ \t]+", " ", texto)       # varios espacios/tabs -> uno solo
    texto = re.sub(r"\n{3,}", "\n\n", texto)    # 3+ saltos de línea seguidos -> 2
    return texto.strip()                         # saca espacios en blanco al principio/final


# --- El "router": decide qué función usar según la extensión del archivo ---

def leer_documento(ruta: Path) -> str:
    """
    Función principal: recibe la ruta de CUALQUIER documento soportado
    y devuelve su texto YA LIMPIO, sin que quien la llama tenga que
    saber de qué formato se trata ni preocuparse por la limpieza.
    """
    extension = ruta.suffix.lower()  # .lower() para que no importe si es .PDF o .pdf

    if extension == ".pdf":
        texto = leer_pdf(ruta)
    elif extension == ".docx":
        texto = leer_word(ruta)
    elif extension == ".xlsx":
        texto = leer_excel(ruta)
    elif extension == ".csv":
        texto = leer_csv(ruta)
    elif extension == ".pptx":
        texto = leer_powerpoint(ruta)
    elif extension == ".html":
        texto = leer_html(ruta)
    elif extension == ".json":
        texto = leer_json(ruta)
    elif extension == ".md":
        texto = reformatear_tablas_markdown(leer_texto_plano(ruta))
    else:
        raise ValueError(f"Formato no soportado: {extension} (archivo: {ruta.name})")

    return limpiar_texto(texto)


# --- Prueba rápida del módulo ---
# Esto SOLO se ejecuta si corrés este archivo directamente
# (python src/agente/parsers.py). Si otro archivo hace "import",
# este bloque no se ejecuta -- es la forma estándar de dejar una
# mini demostración dentro de un módulo sin que moleste al importarlo.
if __name__ == "__main__":
    raiz_del_proyecto = Path(__file__).parent.parent.parent

    documentos_de_prueba = [
        "docs/operacional/FAQ_transacciones_y_limites.md",
        "docs/rh/Politica_de_beneficios_y_vacaciones.pdf",
        "docs/legal_compliance/Terminos_y_condiciones_de_uso.docx",
        "docs/financiero/Tarifas_y_comisiones_del_servicio.xlsx",
        "docs/datos_sistemas/base_clientes_simulada.csv",
        "docs/estrategico/Roadmap_estrategico_2026.pptx",
        "docs/operacional/Manual_tecnico_integracion_api_pagos.html",
        "docs/datos_sistemas/configuracion_endpoints_api.json",
    ]

    for doc_relativo in documentos_de_prueba:
        ruta_doc = raiz_del_proyecto / doc_relativo
        texto = leer_documento(ruta_doc)
        print(f"{ruta_doc.name}: {len(texto)} caracteres extraídos")
